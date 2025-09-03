"""
Document Processing Service for RAG Pipeline
Handles text extraction, chunking, and embedding generation for various file types
"""
import os
import uuid
import json
import mimetypes
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import logging

# Configure logging first
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Text extraction imports
import PyPDF2
import docx
import pandas as pd
import markdown
from pptx import Presentation

# RAG processing imports
import numpy as np
try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_AVAILABLE = True
    logger.info("sentence_transformers successfully imported")
except ImportError as e:
    logger.warning(f"sentence_transformers not available: {e}")
    SENTENCE_TRANSFORMERS_AVAILABLE = False
    SentenceTransformer = None

import nltk
from nltk.tokenize import sent_tokenize
try:
    import tiktoken
    TIKTOKEN_AVAILABLE = True
except ImportError:
    TIKTOKEN_AVAILABLE = False
    tiktoken = None

# Database imports
from models import db, UploadedFile, DocumentChunk, RFPO

class DocumentProcessor:
    """Main document processing service for RAG pipeline"""
    
    def __init__(self, embedding_model_name: str = "all-MiniLM-L6-v2"):
        """
        Initialize the document processor
        
        Args:
            embedding_model_name: Name of the sentence transformer model to use
        """
        self.embedding_model_name = embedding_model_name
        self.embedding_model = None
        self.tokenizer = None
        
        # Download required NLTK data
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')
    
    def _load_embedding_model(self):
        """Lazy load the embedding model to avoid startup delays"""
        if self.embedding_model is None:
            if not SENTENCE_TRANSFORMERS_AVAILABLE:
                raise ImportError("sentence_transformers is not available. Please install it to use RAG functionality.")
            
            logger.info(f"Loading embedding model: {self.embedding_model_name}")
            self.embedding_model = SentenceTransformer(self.embedding_model_name)
            logger.info("Embedding model loaded successfully")
    
    def _load_tokenizer(self):
        """Load tokenizer for token counting"""
        if self.tokenizer is None:
            if TIKTOKEN_AVAILABLE:
                try:
                    self.tokenizer = tiktoken.get_encoding("cl100k_base")
                except Exception as e:
                    logger.warning(f"Could not load tiktoken: {e}. Using character-based estimation.")
                    self.tokenizer = None
            else:
                logger.warning("tiktoken not available. Using character-based estimation.")
                self.tokenizer = None
    
    def extract_text_from_file(self, file_path: str, mime_type: str) -> Tuple[str, Dict]:
        """
        Extract text content from various file formats
        
        Args:
            file_path: Path to the file
            mime_type: MIME type of the file
            
        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        metadata = {
            'extraction_method': 'unknown',
            'page_count': 0,
            'word_count': 0,
            'extraction_timestamp': datetime.utcnow().isoformat()
        }
        
        try:
            if mime_type == 'application/pdf':
                return self._extract_from_pdf(file_path, metadata)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                              'application/msword']:
                return self._extract_from_docx(file_path, metadata)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                              'application/vnd.ms-excel', 'text/csv']:
                return self._extract_from_excel_csv(file_path, metadata)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'application/vnd.ms-powerpoint']:
                return self._extract_from_pptx(file_path, metadata)
            elif mime_type == 'text/markdown':
                return self._extract_from_markdown(file_path, metadata)
            elif mime_type.startswith('text/'):
                return self._extract_from_text(file_path, metadata)
            else:
                # Try to read as text file as fallback
                return self._extract_from_text(file_path, metadata)
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    def _extract_from_pdf(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from PDF files"""
        text_content = ""
        metadata['extraction_method'] = 'PyPDF2'
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            metadata['page_count'] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
        
        metadata['word_count'] = len(text_content.split())
        return text_content.strip(), metadata
    
    def _extract_from_docx(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from DOCX files"""
        metadata['extraction_method'] = 'python-docx'
        
        doc = docx.Document(file_path)
        paragraphs = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text.strip())
        
        text_content = "\n\n".join(paragraphs)
        metadata['paragraph_count'] = len(paragraphs)
        metadata['word_count'] = len(text_content.split())
        
        return text_content, metadata
    
    def _extract_from_excel_csv(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from Excel/CSV files"""
        metadata['extraction_method'] = 'pandas'
        
        try:
            if file_path.lower().endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            # Convert DataFrame to text representation
            text_content = f"Data Summary:\nRows: {len(df)}\nColumns: {len(df.columns)}\n\n"
            text_content += f"Column Names: {', '.join(df.columns)}\n\n"
            
            # Add first few rows as sample
            text_content += "Sample Data:\n"
            text_content += df.head(10).to_string(index=False)
            
            # Add column statistics if numeric columns exist
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text_content += "\n\nNumeric Column Statistics:\n"
                text_content += df[numeric_cols].describe().to_string()
            
            metadata['row_count'] = len(df)
            metadata['column_count'] = len(df.columns)
            metadata['columns'] = df.columns.tolist()
            metadata['word_count'] = len(text_content.split())
            
            return text_content, metadata
            
        except Exception as e:
            logger.error(f"Error processing Excel/CSV file: {e}")
            raise
    
    def _extract_from_pptx(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from PowerPoint files"""
        metadata['extraction_method'] = 'python-pptx'
        
        prs = Presentation(file_path)
        slide_texts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"--- Slide {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
            
            if slide_text.strip() != f"--- Slide {slide_num + 1} ---":
                slide_texts.append(slide_text)
        
        text_content = "\n\n".join(slide_texts)
        metadata['slide_count'] = len(prs.slides)
        metadata['word_count'] = len(text_content.split())
        
        return text_content, metadata
    
    def _extract_from_markdown(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from Markdown files"""
        metadata['extraction_method'] = 'markdown + text'
        
        with open(file_path, 'r', encoding='utf-8') as file:
            markdown_content = file.read()
        
        # Convert markdown to plain text (removes formatting)
        html = markdown.markdown(markdown_content)
        # Simple HTML tag removal
        import re
        text_content = re.sub('<[^<]+?>', '', html)
        
        metadata['word_count'] = len(text_content.split())
        metadata['original_format'] = 'markdown'
        
        return text_content, metadata
    
    def _extract_from_text(self, file_path: str, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from plain text files"""
        metadata['extraction_method'] = 'text'
        
        encodings = ['utf-8', 'latin-1', 'ascii']
        text_content = ""
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text_content = file.read()
                metadata['encoding'] = encoding
                break
            except UnicodeDecodeError:
                continue
        
        if not text_content:
            raise ValueError("Could not decode text file with any supported encoding")
        
        metadata['word_count'] = len(text_content.split())
        return text_content, metadata
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict]:
        """
        Split text into overlapping chunks for better RAG performance
        
        Args:
            text: Input text to chunk
            chunk_size: Target size of each chunk in characters
            overlap: Number of characters to overlap between chunks
            
        Returns:
            List of chunk dictionaries with metadata
        """
        if not text or not text.strip():
            return []
        
        # Split into sentences for better chunk boundaries
        sentences = sent_tokenize(text)
        
        chunks = []
        current_chunk = ""
        current_size = 0
        chunk_index = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            # If adding this sentence would exceed chunk size and we have content
            if current_size + sentence_size > chunk_size and current_chunk:
                # Save current chunk
                chunks.append({
                    'text': current_chunk.strip(),
                    'chunk_index': chunk_index,
                    'size': len(current_chunk.strip()),
                    'sentence_count': len(sent_tokenize(current_chunk.strip()))
                })
                
                # Start new chunk with overlap
                overlap_text = self._get_overlap_text(current_chunk, overlap)
                current_chunk = overlap_text + " " + sentence
                current_size = len(current_chunk)
                chunk_index += 1
            else:
                # Add sentence to current chunk
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
                current_size += sentence_size
        
        # Add final chunk if it has content
        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'chunk_index': chunk_index,
                'size': len(current_chunk.strip()),
                'sentence_count': len(sent_tokenize(current_chunk.strip()))
            })
        
        return chunks
    
    def _get_overlap_text(self, text: str, overlap_size: int) -> str:
        """Get the last part of text for chunk overlap"""
        if len(text) <= overlap_size:
            return text
        
        # Try to break at sentence boundary within overlap
        sentences = sent_tokenize(text)
        overlap_text = ""
        
        for sentence in reversed(sentences):
            if len(overlap_text + sentence) <= overlap_size:
                overlap_text = sentence + " " + overlap_text
            else:
                break
        
        return overlap_text.strip() or text[-overlap_size:]
    
    def generate_embeddings(self, chunks: List[Dict]) -> List[np.ndarray]:
        """
        Generate embeddings for text chunks
        
        Args:
            chunks: List of chunk dictionaries
            
        Returns:
            List of embedding vectors
        """
        if not chunks:
            return []
        
        self._load_embedding_model()
        
        texts = [chunk['text'] for chunk in chunks]
        embeddings = self.embedding_model.encode(texts, show_progress_bar=True)
        
        return embeddings
    
    def process_uploaded_file(self, file_record: UploadedFile) -> bool:
        """
        Complete processing pipeline for an uploaded file
        
        Args:
            file_record: UploadedFile database record
            
        Returns:
            True if processing succeeded, False otherwise
        """
        try:
            logger.info(f"Starting processing for file: {file_record.original_filename}")
            
            # Update status
            file_record.processing_status = 'processing'
            db.session.commit()
            
            # Extract text
            text_content, metadata = self.extract_text_from_file(
                file_record.file_path, 
                file_record.mime_type
            )
            
            if not text_content.strip():
                raise ValueError("No text content extracted from file")
            
            file_record.text_extracted = True
            
            # Chunk text
            chunks = self.chunk_text(text_content)
            
            if not chunks:
                raise ValueError("No chunks generated from text content")
            
            # Generate embeddings
            embeddings = self.generate_embeddings(chunks)
            
            # Save chunks to database
            for i, (chunk_data, embedding) in enumerate(zip(chunks, embeddings)):
                chunk = DocumentChunk(
                    chunk_id=str(uuid.uuid4()),
                    text_content=chunk_data['text'],
                    chunk_index=chunk_data['chunk_index'],
                    chunk_size=chunk_data['size'],
                    file_id=file_record.id,
                    embedding_model=self.embedding_model_name
                )
                
                # Store embedding and metadata
                chunk.set_embedding(embedding)
                chunk.set_metadata({
                    'sentence_count': chunk_data['sentence_count'],
                    'extraction_metadata': metadata
                })
                
                db.session.add(chunk)
            
            # Update file record
            file_record.embeddings_created = True
            file_record.chunk_count = len(chunks)
            file_record.processing_status = 'completed'
            file_record.processed_at = datetime.utcnow()
            
            db.session.commit()
            
            logger.info(f"Successfully processed file: {file_record.original_filename} "
                       f"({len(chunks)} chunks created)")
            
            return True
            
        except Exception as e:
            logger.error(f"Error processing file {file_record.original_filename}: {str(e)}")
            
            # Update error status
            file_record.processing_status = 'failed'
            file_record.processing_error = str(e)
            db.session.commit()
            
            return False
    
    def search_similar_chunks(self, query: str, rfpo_id: int, top_k: int = 5) -> List[Dict]:
        """
        Search for similar text chunks using vector similarity
        
        Args:
            query: Search query
            rfpo_id: RFPO ID to limit search scope
            top_k: Number of top results to return
            
        Returns:
            List of similar chunks with metadata
        """
        self._load_embedding_model()
        
        # Generate query embedding
        query_embedding = self.embedding_model.encode([query])[0]
        
        # Get all chunks for the RFPO with explicit join and fresh session
        chunks_query = db.session.query(DocumentChunk, UploadedFile).join(
            UploadedFile, DocumentChunk.file_id == UploadedFile.id
        ).filter(
            UploadedFile.rfpo_id == rfpo_id,
            DocumentChunk.embedding_vector.isnot(None)
        )
        
        chunk_file_pairs = chunks_query.all()
        
        if not chunk_file_pairs:
            return []
        
        # Calculate similarities
        similarities = []
        for chunk, file_record in chunk_file_pairs:
            try:
                embedding_data = chunk.get_embedding()
                if embedding_data:
                    chunk_embedding = np.array(embedding_data)
                    similarity = np.dot(query_embedding, chunk_embedding) / (
                        np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                    )
                    similarities.append((chunk, file_record, similarity))
            except Exception as e:
                logger.warning(f"Error calculating similarity for chunk {chunk.id}: {str(e)}")
                continue
        
        # Sort by similarity and return top k
        similarities.sort(key=lambda x: x[2], reverse=True)
        
        results = []
        for chunk, file_record, similarity in similarities[:top_k]:
            result = chunk.to_dict()
            result['similarity_score'] = float(similarity)
            result['file_name'] = file_record.original_filename
            results.append(result)
        
        return results

# Global instance
document_processor = DocumentProcessor()
